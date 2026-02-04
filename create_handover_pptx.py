#!/usr/bin/env python3
"""
Génération du PowerPoint de passation ArchiMeuble
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Couleurs du thème
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)       # Fond sombre
ACCENT_GOLD = RGBColor(0xC9, 0x96, 0x3B)    # Or/doré
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MEDIUM_GRAY = RGBColor(0x99, 0x99, 0x99)
SOFT_WHITE = RGBColor(0xF0, 0xF0, 0xF0)
CARD_BG = RGBColor(0x2A, 0x2A, 0x42)        # Fond carte
SECTION_BG = RGBColor(0x22, 0x22, 0x3A)     # Fond section

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_dark_background(slide):
    """Ajoute un fond sombre à la slide"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK_BG


def add_shape_with_fill(slide, left, top, width, height, color):
    """Ajoute un rectangle arrondi avec couleur de fond"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    # Ajuster le rayon d'arrondi
    shape.adjustments[0] = 0.05
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    """Ajoute une zone de texte"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_rich_text_box(slide, left, top, width, height, lines):
    """
    Ajoute une zone de texte avec plusieurs lignes formatées.
    lines = [(text, font_size, color, bold, alignment), ...]
    """
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, (text, font_size, color, bold, alignment) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = alignment
        p.space_after = Pt(6)

    return txBox


def add_gold_line(slide, left, top, width):
    """Ajoute une ligne dorée décorative"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_GOLD
    shape.line.fill.background()
    return shape


def add_credential_card(slide, left, top, width, height, title, items):
    """
    Ajoute une carte d'identifiants.
    items = [(label, value), ...]
    """
    card = add_shape_with_fill(slide, left, top, width, height, CARD_BG)

    # Titre de la carte
    add_text_box(slide, left + Inches(0.3), top + Inches(0.15), width - Inches(0.6), Inches(0.4),
                 title, font_size=16, color=ACCENT_GOLD, bold=True)

    # Ligne séparatrice
    add_gold_line(slide, left + Inches(0.3), top + Inches(0.55), width - Inches(0.6))

    # Items
    y_offset = top + Inches(0.7)
    for label, value in items:
        add_text_box(slide, left + Inches(0.3), y_offset, width - Inches(0.6), Inches(0.25),
                     label, font_size=11, color=MEDIUM_GRAY, bold=False)
        y_offset += Inches(0.22)
        add_text_box(slide, left + Inches(0.3), y_offset, width - Inches(0.6), Inches(0.3),
                     value, font_size=13, color=WHITE, bold=True)
        y_offset += Inches(0.35)

    return card


# =============================================================================
# SLIDE 1 - Page de garde
# =============================================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Layout vide
add_dark_background(slide1)

# Ligne dorée en haut
add_gold_line(slide1, Inches(0), Inches(0), SLIDE_W)

# Titre principal
add_rich_text_box(slide1, Inches(1), Inches(1.8), Inches(11.3), Inches(4), [
    ("ARCHIMEUBLE", 48, ACCENT_GOLD, True, PP_ALIGN.CENTER),
    ("", 12, WHITE, False, PP_ALIGN.CENTER),
    ("Document de Passation", 32, WHITE, True, PP_ALIGN.CENTER),
    ("", 12, WHITE, False, PP_ALIGN.CENTER),
    ("Guide complet des services et accès", 18, LIGHT_GRAY, False, PP_ALIGN.CENTER),
    ("", 24, WHITE, False, PP_ALIGN.CENTER),
    ("Janvier 2026", 14, MEDIUM_GRAY, False, PP_ALIGN.CENTER),
])

# Ligne dorée en bas
add_gold_line(slide1, Inches(2), Inches(6.2), Inches(9.3))

# =============================================================================
# SLIDE 2 - Sommaire
# =============================================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide2)

add_text_box(slide2, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             "Sommaire", font_size=36, color=ACCENT_GOLD, bold=True)
add_gold_line(slide2, Inches(0.8), Inches(1.2), Inches(4))

sommaire_items = [
    ("1.", "Le Projet ArchiMeuble", "Qu'est-ce que c'est et comment ça fonctionne"),
    ("2.", "GitHub", "Le code source du site"),
    ("3.", "Vercel", "L'hébergement du site web (ce que voient les visiteurs)"),
    ("4.", "Railway", "L'hébergement du serveur (le moteur invisible)"),
    ("5.", "OVH", "Le nom de domaine archimeuble.com"),
    ("6.", "Stripe", "Les paiements en ligne"),
    ("7.", "Brevo", "L'envoi d'emails automatiques"),
    ("8.", "Calendly", "La prise de rendez-vous en ligne"),
    ("9.", "Récapitulatif des accès", "Tous les identifiants en un coup d'oeil"),
]

y = Inches(1.6)
for num, title, desc in sommaire_items:
    add_text_box(slide2, Inches(1.2), y, Inches(0.5), Inches(0.35),
                 num, font_size=16, color=ACCENT_GOLD, bold=True)
    add_text_box(slide2, Inches(1.7), y, Inches(4), Inches(0.35),
                 title, font_size=16, color=WHITE, bold=True)
    add_text_box(slide2, Inches(5.8), y, Inches(6), Inches(0.35),
                 desc, font_size=13, color=LIGHT_GRAY, bold=False)
    y += Inches(0.55)

# =============================================================================
# SLIDE 3 - Le Projet ArchiMeuble
# =============================================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide3)

add_text_box(slide3, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             "1. Le Projet ArchiMeuble", font_size=36, color=ACCENT_GOLD, bold=True)
add_gold_line(slide3, Inches(0.8), Inches(1.2), Inches(5))

# Carte "C'est quoi ?"
card1 = add_shape_with_fill(slide3, Inches(0.8), Inches(1.6), Inches(5.5), Inches(2.5), CARD_BG)
add_rich_text_box(slide3, Inches(1.1), Inches(1.7), Inches(5), Inches(2.3), [
    ("C'est quoi ?", 20, ACCENT_GOLD, True, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("ArchiMeuble est un site web qui permet aux clients", 14, WHITE, False, PP_ALIGN.LEFT),
    ("de configurer des meubles sur mesure en 3D,", 14, WHITE, False, PP_ALIGN.LEFT),
    ("de visualiser le résultat, et de passer commande", 14, WHITE, False, PP_ALIGN.LEFT),
    ("directement en ligne avec paiement sécurisé.", 14, WHITE, False, PP_ALIGN.LEFT),
])

# Carte "Comment ça marche ?"
card2 = add_shape_with_fill(slide3, Inches(7), Inches(1.6), Inches(5.5), Inches(2.5), CARD_BG)
add_rich_text_box(slide3, Inches(7.3), Inches(1.7), Inches(5), Inches(2.3), [
    ("Comment ça marche ?", 20, ACCENT_GOLD, True, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("1. Le client choisit un modèle de meuble", 14, WHITE, False, PP_ALIGN.LEFT),
    ("2. Il personnalise les dimensions, matériaux, couleurs", 14, WHITE, False, PP_ALIGN.LEFT),
    ("3. Il voit le rendu 3D en temps réel", 14, WHITE, False, PP_ALIGN.LEFT),
    ("4. Il passe commande et paie en ligne", 14, WHITE, False, PP_ALIGN.LEFT),
])

# Carte "Architecture technique simplifiée"
card3 = add_shape_with_fill(slide3, Inches(0.8), Inches(4.5), Inches(11.7), Inches(2.5), CARD_BG)
add_rich_text_box(slide3, Inches(1.1), Inches(4.6), Inches(11.2), Inches(2.3), [
    ("Le site est composé de 2 parties", 20, ACCENT_GOLD, True, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("Le Frontend (la vitrine)  -  C'est ce que le visiteur voit : les pages, les boutons, le configurateur 3D.", 14, WHITE, False, PP_ALIGN.LEFT),
    ("     Hébergé sur Vercel  |  Adresses : staging.archimeuble.com / dev.archimeuble.com", 13, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("Le Backend (le moteur)  -  C'est la partie invisible : gestion des commandes, paiements, comptes clients, emails.", 14, WHITE, False, PP_ALIGN.LEFT),
    ("     Hébergé sur Railway  |  Adresses : api-staging.archimeuble.com / api-dev.archimeuble.com", 13, LIGHT_GRAY, False, PP_ALIGN.LEFT),
])


# =============================================================================
# SLIDE 4 - GitHub
# =============================================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide4)

add_text_box(slide4, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             "2. GitHub - Le code source", font_size=36, color=ACCENT_GOLD, bold=True)
add_gold_line(slide4, Inches(0.8), Inches(1.2), Inches(5))

# Explication
card_gh1 = add_shape_with_fill(slide4, Inches(0.8), Inches(1.6), Inches(11.7), Inches(1.8), CARD_BG)
add_rich_text_box(slide4, Inches(1.1), Inches(1.7), Inches(11.2), Inches(1.6), [
    ("C'est quoi GitHub ?", 20, ACCENT_GOLD, True, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("GitHub est comme un coffre-fort pour le code du site. Tout le code source y est stocké en sécurité.", 14, WHITE, False, PP_ALIGN.LEFT),
    ("Si quelqu'un doit modifier le site, c'est ici qu'il trouvera tout le code. C'est aussi ici que Vercel et Railway", 14, WHITE, False, PP_ALIGN.LEFT),
    ("vont chercher le code pour le publier en ligne automatiquement.", 14, WHITE, False, PP_ALIGN.LEFT),
])

# Les deux dépôts
card_gh2 = add_shape_with_fill(slide4, Inches(0.8), Inches(3.8), Inches(5.5), Inches(3.2), CARD_BG)
add_rich_text_box(slide4, Inches(1.1), Inches(3.9), Inches(5), Inches(3), [
    ("Dépôt Frontend (le site visible)", 18, ACCENT_GOLD, True, PP_ALIGN.LEFT),
    ("", 6, WHITE, False, PP_ALIGN.LEFT),
    ("Organisation : archimeuble", 13, WHITE, False, PP_ALIGN.LEFT),
    ("Nom du dépôt : front", 13, WHITE, False, PP_ALIGN.LEFT),
    ("", 6, WHITE, False, PP_ALIGN.LEFT),
    ("Branches (versions) :", 13, ACCENT_GOLD, True, PP_ALIGN.LEFT),
    ("  dev  =  version de développement (tests)", 13, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("  staging  =  version de pré-production (validation)", 13, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("  main  =  version finale (production)", 13, LIGHT_GRAY, False, PP_ALIGN.LEFT),
])

card_gh3 = add_shape_with_fill(slide4, Inches(7), Inches(3.8), Inches(5.5), Inches(3.2), CARD_BG)
add_rich_text_box(slide4, Inches(7.3), Inches(3.9), Inches(5), Inches(3), [
    ("Dépôt Backend (le moteur)", 18, ACCENT_GOLD, True, PP_ALIGN.LEFT),
    ("", 6, WHITE, False, PP_ALIGN.LEFT),
    ("Organisation : archimeuble", 13, WHITE, False, PP_ALIGN.LEFT),
    ("Nom du dépôt : back", 13, WHITE, False, PP_ALIGN.LEFT),
    ("", 6, WHITE, False, PP_ALIGN.LEFT),
    ("Branches (versions) :", 13, ACCENT_GOLD, True, PP_ALIGN.LEFT),
    ("  dev  =  version de développement (tests)", 13, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("  staging  =  version de pré-production (validation)", 13, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("  main  =  version finale (production)", 13, LIGHT_GRAY, False, PP_ALIGN.LEFT),
])


# =============================================================================
# SLIDE 5 - Vercel
# =============================================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide5)

add_text_box(slide5, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             "3. Vercel - L'hébergement du site web", font_size=36, color=ACCENT_GOLD, bold=True)
add_gold_line(slide5, Inches(0.8), Inches(1.2), Inches(6))

card_v1 = add_shape_with_fill(slide5, Inches(0.8), Inches(1.6), Inches(11.7), Inches(2.0), CARD_BG)
add_rich_text_box(slide5, Inches(1.1), Inches(1.7), Inches(11.2), Inches(1.8), [
    ("C'est quoi Vercel ?", 20, ACCENT_GOLD, True, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("Vercel est le service qui rend le site web accessible aux visiteurs. Quand quelqu'un tape", 14, WHITE, False, PP_ALIGN.LEFT),
    ("archimeuble.com dans son navigateur, c'est Vercel qui affiche les pages.", 14, WHITE, False, PP_ALIGN.LEFT),
    ("Vercel se connecte automatiquement à GitHub : dès qu'une modification est faite dans le code, le site se met à jour.", 14, WHITE, False, PP_ALIGN.LEFT),
])

card_v2 = add_shape_with_fill(slide5, Inches(0.8), Inches(4.0), Inches(5.5), Inches(3.0), CARD_BG)
add_rich_text_box(slide5, Inches(1.1), Inches(4.1), Inches(5), Inches(2.8), [
    ("Les environnements", 18, ACCENT_GOLD, True, PP_ALIGN.LEFT),
    ("", 6, WHITE, False, PP_ALIGN.LEFT),
    ("Production (branche staging)", 14, WHITE, True, PP_ALIGN.LEFT),
    ("  staging.archimeuble.com", 13, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("  C'est la version visible par les clients", 12, MEDIUM_GRAY, False, PP_ALIGN.LEFT),
    ("", 6, WHITE, False, PP_ALIGN.LEFT),
    ("Preview (branche dev)", 14, WHITE, True, PP_ALIGN.LEFT),
    ("  dev.archimeuble.com", 13, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("  C'est la version de test pour les développeurs", 12, MEDIUM_GRAY, False, PP_ALIGN.LEFT),
])

add_credential_card(slide5, Inches(7), Inches(4.0), Inches(5.5), Inches(2.2),
    "Connexion Vercel", [
        ("Méthode", "Se connecter avec GitHub"),
        ("Compte GitHub", "Organisation archimeuble"),
    ])


# =============================================================================
# SLIDE 6 - Railway
# =============================================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide6)

add_text_box(slide6, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             "4. Railway - L'hébergement du serveur", font_size=36, color=ACCENT_GOLD, bold=True)
add_gold_line(slide6, Inches(0.8), Inches(1.2), Inches(6))

card_r1 = add_shape_with_fill(slide6, Inches(0.8), Inches(1.6), Inches(11.7), Inches(2.0), CARD_BG)
add_rich_text_box(slide6, Inches(1.1), Inches(1.7), Inches(11.2), Inches(1.8), [
    ("C'est quoi Railway ?", 20, ACCENT_GOLD, True, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("Railway héberge le serveur (backend) du site. C'est la partie invisible qui gère :", 14, WHITE, False, PP_ALIGN.LEFT),
    ("les comptes clients, les commandes, les paiements, l'envoi d'emails, et la base de données.", 14, WHITE, False, PP_ALIGN.LEFT),
    ("Comme Vercel, il se connecte à GitHub et se met à jour automatiquement.", 14, WHITE, False, PP_ALIGN.LEFT),
])

card_r2 = add_shape_with_fill(slide6, Inches(0.8), Inches(4.0), Inches(5.5), Inches(3.0), CARD_BG)
add_rich_text_box(slide6, Inches(1.1), Inches(4.1), Inches(5), Inches(2.8), [
    ("Les environnements", 18, ACCENT_GOLD, True, PP_ALIGN.LEFT),
    ("", 6, WHITE, False, PP_ALIGN.LEFT),
    ("Staging (pré-production)", 14, WHITE, True, PP_ALIGN.LEFT),
    ("  api-staging.archimeuble.com", 13, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("  Le serveur utilisé par staging.archimeuble.com", 12, MEDIUM_GRAY, False, PP_ALIGN.LEFT),
    ("", 6, WHITE, False, PP_ALIGN.LEFT),
    ("Dev (développement)", 14, WHITE, True, PP_ALIGN.LEFT),
    ("  api-dev.archimeuble.com", 13, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("  Le serveur utilisé par dev.archimeuble.com", 12, MEDIUM_GRAY, False, PP_ALIGN.LEFT),
])

add_credential_card(slide6, Inches(7), Inches(4.0), Inches(5.5), Inches(2.2),
    "Connexion Railway", [
        ("Méthode", "Se connecter avec GitHub"),
        ("Compte GitHub", "Organisation archimeuble"),
    ])


# =============================================================================
# SLIDE 7 - OVH
# =============================================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide7)

add_text_box(slide7, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             "5. OVH - Le nom de domaine", font_size=36, color=ACCENT_GOLD, bold=True)
add_gold_line(slide7, Inches(0.8), Inches(1.2), Inches(5))

card_o1 = add_shape_with_fill(slide7, Inches(0.8), Inches(1.6), Inches(11.7), Inches(2.2), CARD_BG)
add_rich_text_box(slide7, Inches(1.1), Inches(1.7), Inches(11.2), Inches(2.0), [
    ("C'est quoi OVH ?", 20, ACCENT_GOLD, True, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("OVH est le service où le nom de domaine archimeuble.com est acheté et géré.", 14, WHITE, False, PP_ALIGN.LEFT),
    ("C'est ici qu'on configure les adresses du site (DNS) pour que les visiteurs soient dirigés", 14, WHITE, False, PP_ALIGN.LEFT),
    ("vers les bons serveurs (Vercel pour le site, Railway pour le serveur).", 14, WHITE, False, PP_ALIGN.LEFT),
])

card_o2 = add_shape_with_fill(slide7, Inches(0.8), Inches(4.2), Inches(11.7), Inches(2.8), CARD_BG)
add_rich_text_box(slide7, Inches(1.1), Inches(4.3), Inches(11.2), Inches(2.6), [
    ("Configuration DNS actuelle", 18, ACCENT_GOLD, True, PP_ALIGN.LEFT),
    ("", 6, WHITE, False, PP_ALIGN.LEFT),
    ("staging.archimeuble.com  -->  Vercel (site de pré-production)", 13, WHITE, False, PP_ALIGN.LEFT),
    ("dev.archimeuble.com  -->  Vercel (site de développement)", 13, WHITE, False, PP_ALIGN.LEFT),
    ("api-staging.archimeuble.com  -->  Railway (serveur de pré-production)", 13, WHITE, False, PP_ALIGN.LEFT),
    ("api-dev.archimeuble.com  -->  Railway (serveur de développement)", 13, WHITE, False, PP_ALIGN.LEFT),
    ("", 6, WHITE, False, PP_ALIGN.LEFT),
    ("Note : L'accès OVH est géré par le propriétaire du domaine.", 12, MEDIUM_GRAY, False, PP_ALIGN.LEFT),
])


# =============================================================================
# SLIDE 8 - Stripe
# =============================================================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide8)

add_text_box(slide8, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             "6. Stripe - Les paiements en ligne", font_size=36, color=ACCENT_GOLD, bold=True)
add_gold_line(slide8, Inches(0.8), Inches(1.2), Inches(6))

card_s1 = add_shape_with_fill(slide8, Inches(0.8), Inches(1.6), Inches(11.7), Inches(2.2), CARD_BG)
add_rich_text_box(slide8, Inches(1.1), Inches(1.7), Inches(11.2), Inches(2.0), [
    ("C'est quoi Stripe ?", 20, ACCENT_GOLD, True, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("Stripe est le service qui gère les paiements par carte bancaire sur le site.", 14, WHITE, False, PP_ALIGN.LEFT),
    ("Quand un client paie sa commande, c'est Stripe qui traite le paiement de manière sécurisée.", 14, WHITE, False, PP_ALIGN.LEFT),
    ("L'argent des ventes arrive sur le compte Stripe, puis peut être transféré vers un compte bancaire.", 14, WHITE, False, PP_ALIGN.LEFT),
])

card_s2 = add_shape_with_fill(slide8, Inches(0.8), Inches(4.2), Inches(5.5), Inches(2.5), CARD_BG)
add_rich_text_box(slide8, Inches(1.1), Inches(4.3), Inches(5), Inches(2.3), [
    ("Mode Test vs Production", 18, ACCENT_GOLD, True, PP_ALIGN.LEFT),
    ("", 6, WHITE, False, PP_ALIGN.LEFT),
    ("Actuellement en mode Test", 14, WHITE, True, PP_ALIGN.LEFT),
    ("Les paiements ne sont pas réels.", 13, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("", 6, WHITE, False, PP_ALIGN.LEFT),
    ("Pour accepter de vrais paiements :", 14, WHITE, True, PP_ALIGN.LEFT),
    ("Activer le mode Production dans Stripe", 13, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("et remplacer les clés de test par celles de prod.", 13, LIGHT_GRAY, False, PP_ALIGN.LEFT),
])

add_credential_card(slide8, Inches(7), Inches(4.2), Inches(5.5), Inches(2.5),
    "Connexion Stripe", [
        ("Email", "pro.archimeuble@gmail.com"),
        ("Mot de passe", "archimeuble2025!"),
    ])


# =============================================================================
# SLIDE 9 - Brevo
# =============================================================================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide9)

add_text_box(slide9, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             "7. Brevo - L'envoi d'emails", font_size=36, color=ACCENT_GOLD, bold=True)
add_gold_line(slide9, Inches(0.8), Inches(1.2), Inches(5))

card_b1 = add_shape_with_fill(slide9, Inches(0.8), Inches(1.6), Inches(11.7), Inches(2.0), CARD_BG)
add_rich_text_box(slide9, Inches(1.1), Inches(1.7), Inches(11.2), Inches(1.8), [
    ("C'est quoi Brevo ?", 20, ACCENT_GOLD, True, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("Brevo (anciennement Sendinblue) est le service d'envoi d'emails automatiques.", 14, WHITE, False, PP_ALIGN.LEFT),
    ("Il envoie les emails de confirmation de commande, de réinitialisation de mot de passe,", 14, WHITE, False, PP_ALIGN.LEFT),
    ("et de notification aux administrateurs quand un client passe une commande.", 14, WHITE, False, PP_ALIGN.LEFT),
])

card_b2 = add_shape_with_fill(slide9, Inches(0.8), Inches(4.0), Inches(5.5), Inches(2.5), CARD_BG)
add_rich_text_box(slide9, Inches(1.1), Inches(4.1), Inches(5), Inches(2.3), [
    ("Types d'emails envoyés", 18, ACCENT_GOLD, True, PP_ALIGN.LEFT),
    ("", 6, WHITE, False, PP_ALIGN.LEFT),
    ("Confirmation de commande au client", 13, WHITE, False, PP_ALIGN.LEFT),
    ("Notification de nouvelle commande à l'admin", 13, WHITE, False, PP_ALIGN.LEFT),
    ("Réinitialisation de mot de passe", 13, WHITE, False, PP_ALIGN.LEFT),
    ("Confirmation de paiement", 13, WHITE, False, PP_ALIGN.LEFT),
    ("Notification de nouvelle configuration", 13, WHITE, False, PP_ALIGN.LEFT),
])

add_credential_card(slide9, Inches(7), Inches(4.0), Inches(5.5), Inches(2.2),
    "Connexion Brevo", [
        ("Email", "pro.archimeuble@gmail.com"),
        ("Mot de passe", "Voir gestionnaire de mots de passe"),
    ])


# =============================================================================
# SLIDE 10 - Calendly
# =============================================================================
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide10)

add_text_box(slide10, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             "8. Calendly - La prise de rendez-vous", font_size=36, color=ACCENT_GOLD, bold=True)
add_gold_line(slide10, Inches(0.8), Inches(1.2), Inches(6))

card_c1 = add_shape_with_fill(slide10, Inches(0.8), Inches(1.6), Inches(11.7), Inches(2.0), CARD_BG)
add_rich_text_box(slide10, Inches(1.1), Inches(1.7), Inches(11.2), Inches(1.8), [
    ("C'est quoi Calendly ?", 20, ACCENT_GOLD, True, PP_ALIGN.LEFT),
    ("", 8, WHITE, False, PP_ALIGN.LEFT),
    ("Calendly permet aux clients de prendre rendez-vous directement depuis le site.", 14, WHITE, False, PP_ALIGN.LEFT),
    ("Les clients choisissent un créneau disponible et le rendez-vous est automatiquement", 14, WHITE, False, PP_ALIGN.LEFT),
    ("ajouté au calendrier. Aucune intervention manuelle n'est nécessaire.", 14, WHITE, False, PP_ALIGN.LEFT),
])

card_c2 = add_shape_with_fill(slide10, Inches(0.8), Inches(4.0), Inches(5.5), Inches(2.5), CARD_BG)
add_rich_text_box(slide10, Inches(1.1), Inches(4.1), Inches(5), Inches(2.3), [
    ("Types de rendez-vous", 18, ACCENT_GOLD, True, PP_ALIGN.LEFT),
    ("", 6, WHITE, False, PP_ALIGN.LEFT),
    ("Consultation initiale", 14, WHITE, True, PP_ALIGN.LEFT),
    ("  Premier contact avec le client", 12, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("", 6, WHITE, False, PP_ALIGN.LEFT),
    ("Suivi de projet", 14, WHITE, True, PP_ALIGN.LEFT),
    ("  Point d'avancement sur une commande", 12, LIGHT_GRAY, False, PP_ALIGN.LEFT),
])

add_credential_card(slide10, Inches(7), Inches(4.0), Inches(5.5), Inches(2.5),
    "Connexion Calendly", [
        ("Email", "pro.archimeuble@gmail.com"),
        ("Mot de passe", "archimeuble2025!#"),
    ])


# =============================================================================
# SLIDE 11 - Récapitulatif des accès
# =============================================================================
slide11 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide11)

add_text_box(slide11, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             "9. Récapitulatif de tous les accès", font_size=36, color=ACCENT_GOLD, bold=True)
add_gold_line(slide11, Inches(0.8), Inches(1.2), Inches(6))

# Ligne 1 - GitHub, Vercel, Railway
add_credential_card(slide11, Inches(0.5), Inches(1.6), Inches(3.9), Inches(2.3),
    "GitHub", [
        ("Méthode", "Connexion GitHub"),
        ("Organisation", "archimeuble"),
    ])

add_credential_card(slide11, Inches(4.7), Inches(1.6), Inches(3.9), Inches(2.3),
    "Vercel", [
        ("Méthode", "Connexion via GitHub"),
        ("Projet", "archimeuble-front"),
    ])

add_credential_card(slide11, Inches(8.9), Inches(1.6), Inches(3.9), Inches(2.3),
    "Railway", [
        ("Méthode", "Connexion via GitHub"),
        ("Compte", "pro.archimeuble@gmail.com"),
    ])

# Ligne 2 - Stripe, Brevo, Calendly
add_credential_card(slide11, Inches(0.5), Inches(4.3), Inches(3.9), Inches(2.8),
    "Stripe", [
        ("Email", "pro.archimeuble@gmail.com"),
        ("Mot de passe", "archimeuble2025!"),
        ("Mode", "Test (à activer en prod)"),
    ])

add_credential_card(slide11, Inches(4.7), Inches(4.3), Inches(3.9), Inches(2.3),
    "Brevo", [
        ("Email", "pro.archimeuble@gmail.com"),
        ("Mot de passe", "Voir gestionnaire"),
    ])

add_credential_card(slide11, Inches(8.9), Inches(4.3), Inches(3.9), Inches(2.8),
    "Calendly", [
        ("Email", "pro.archimeuble@gmail.com"),
        ("Mot de passe", "archimeuble2025!#"),
    ])


# =============================================================================
# SLIDE 12 - Workflow de déploiement
# =============================================================================
slide12 = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_background(slide12)

add_text_box(slide12, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7),
             "Comment publier une modification", font_size=36, color=ACCENT_GOLD, bold=True)
add_gold_line(slide12, Inches(0.8), Inches(1.2), Inches(6))

card_w1 = add_shape_with_fill(slide12, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.4), CARD_BG)
add_rich_text_box(slide12, Inches(1.1), Inches(1.7), Inches(11.2), Inches(5.2), [
    ("Le processus de mise en ligne (pour le développeur)", 20, ACCENT_GOLD, True, PP_ALIGN.LEFT),
    ("", 10, WHITE, False, PP_ALIGN.LEFT),
    ("Etape 1 : Développer sur la branche dev", 16, WHITE, True, PP_ALIGN.LEFT),
    ("Le développeur modifie le code sur la branche dev. Le site dev.archimeuble.com se met à jour", 14, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("automatiquement pour tester les modifications.", 14, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("", 10, WHITE, False, PP_ALIGN.LEFT),
    ("Etape 2 : Valider sur staging", 16, WHITE, True, PP_ALIGN.LEFT),
    ("Une fois les tests terminés, on fusionne (merge) le code de dev vers staging.", 14, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("Le site staging.archimeuble.com se met à jour automatiquement pour validation finale.", 14, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("", 10, WHITE, False, PP_ALIGN.LEFT),
    ("Etape 3 : Mettre en production", 16, WHITE, True, PP_ALIGN.LEFT),
    ("Après validation, on fusionne staging vers main. Le site archimeuble.com est mis à jour.", 14, LIGHT_GRAY, False, PP_ALIGN.LEFT),
    ("", 10, WHITE, False, PP_ALIGN.LEFT),
    ("Résumé :   dev  -->  staging  -->  main (production)", 16, ACCENT_GOLD, True, PP_ALIGN.CENTER),
])


# =============================================================================
# SAUVEGARDE
# =============================================================================
output_path = "/Users/collins/Desktop/ArchiMeuble_Passation.pptx"
prs.save(output_path)
print(f"PowerPoint sauvegardé : {output_path}")
